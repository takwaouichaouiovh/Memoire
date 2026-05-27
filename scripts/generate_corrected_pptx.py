"""Generate a corrected version of the soutenance PPTX aligned with the real codebase."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ─── Brand palette ──────────────────────────────────────────────────────────
BG = RGBColor(0x0F, 0x17, 0x2A)          # slate-900
PANEL = RGBColor(0x1E, 0x29, 0x3B)       # slate-800
VIOLET = RGBColor(0x7C, 0x3A, 0xED)      # accent
CYAN = RGBColor(0x22, 0xD3, 0xEE)        # accent 2
OVH = RGBColor(0x12, 0x3F, 0x6D)         # OVH blue
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_slide(prs: Presentation, bg: RGBColor = BG):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.shadow.inherit = False
    return slide


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.adjustments[0] = 0.08
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def add_footer(slide, prs, label):
    add_text(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3),
             label, size=9, color=MUTED)


def add_header_bar(slide, prs):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = VIOLET
    bar.line.fill.background()


def title_slide(slide, prs, eyebrow, title, subtitle):
    add_header_bar(slide, prs)
    add_text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.4),
             eyebrow, size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.9),
             title, size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.65), Inches(12), Inches(0.5),
                 subtitle, size=14, color=MUTED)


FOOTER = "PO.ai — Copilot IA pour Product Owners · Aouichaoui Takwa · M2 MIAS Centrale Lille · OVHcloud IAWF · 2026"


def build():
    prs = new_prs()

    # ─── Slide 1 — Cover ─────────────────────────────────────────────────
    s = add_slide(prs)
    add_header_bar(s, prs)
    add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
             "M2 MIAS — ÉCOLE CENTRALE DE LILLE — 2026", size=14, bold=True, color=CYAN)
    add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(1.2),
             "PO.ai", size=72, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(3.5), Inches(12), Inches(0.6),
             "Un copilote IA pour Product Owners", size=24, color=WHITE)
    add_text(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.5),
             "Assistant RAG  ·  7 algorithmes de priorisation  ·  Agents IA (ReAct + LangGraph)",
             size=16, color=MUTED)
    add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
             "Aouichaoui Takwa  |  Stage OVHcloud — IA Web Factory  |  Juillet 2026",
             size=13, color=MUTED)

    # ─── Slide 2 — Problem ───────────────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "PROBLÉMATIQUE", "Le Product Owner est saturé",
                "Trois douleurs identifiées chez les PO de l'IA Web Factory (OVHcloud)")
    cards = [
        ("📄  Friction documentaire",
         "3–5 h/semaine perdues à chercher dans Confluence et Jira. Aucun moyen d'interroger les docs en langage naturel."),
        ("🎯  Priorisation subjective",
         "Décisions à l'intuition ou au stakeholder le plus insistant. Pas d'outil structuré chez OVH pour la priorisation multi-critères."),
        ("🧠  Biais cognitifs",
         "Biais de disponibilité, confirmation, effet de halo (Kahneman 2011). Aucun outil existant ne les compense."),
    ]
    for i, (h, b) in enumerate(cards):
        left = Inches(0.6 + i * 4.15)
        add_rect(s, left, Inches(2.6), Inches(3.9), Inches(3.4), PANEL)
        add_text(s, left + Inches(0.25), Inches(2.85), Inches(3.4), Inches(0.6),
                 h, size=16, bold=True, color=CYAN)
        add_text(s, left + Inches(0.25), Inches(3.6), Inches(3.4), Inches(2.2),
                 b, size=12, color=WHITE)
    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
             "Question de recherche : comment l'IA peut-elle augmenter le PO sans le remplacer, "
             "en combinant RAG documentaire et priorisation multi-algorithmes ?",
             size=12, bold=True, color=CYAN)
    add_footer(s, prs, FOOTER)

    # ─── Slide 3 — Vision ───────────────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "VISION", "PO.ai — Un copilote, pas un pilote automatique",
                "Trois piliers complémentaires")
    pillars = [
        ("📚  SAVOIR", "RAG documentaire",
         "Interroger Confluence, PDFs et guidelines en langage naturel. Sources citées à chaque réponse."),
        ("⚖️  DÉCIDER", "7 algorithmes",
         "RICE · WSJF · ICE · Kano · Value/Effort · AI Blend · ML Hybrid. Comparaison cross-algos en 1 clic."),
        ("🚀  AGIR", "Agents IA",
         "Tool-calling ReAct · Epic Grooming LangGraph · Sprint Planner knapsack · Retro Analyzer."),
    ]
    for i, (h, sub, body) in enumerate(pillars):
        left = Inches(0.6 + i * 4.15)
        add_rect(s, left, Inches(2.6), Inches(3.9), Inches(3.6), PANEL)
        add_text(s, left + Inches(0.25), Inches(2.85), Inches(3.4), Inches(0.5),
                 h, size=18, bold=True, color=VIOLET)
        add_text(s, left + Inches(0.25), Inches(3.4), Inches(3.4), Inches(0.5),
                 sub, size=14, bold=True, color=WHITE)
        add_text(s, left + Inches(0.25), Inches(4.0), Inches(3.4), Inches(2.0),
                 body, size=12, color=MUTED)
    add_footer(s, prs, FOOTER)

    # ─── Slide 4 — Part 1 ───────────────────────────────────────────────
    s = add_slide(prs, BG)
    add_text(s, Inches(0.6), Inches(2.8), Inches(12), Inches(0.5),
             "PARTIE 1", size=18, bold=True, color=CYAN)
    add_text(s, Inches(0.6), Inches(3.3), Inches(12), Inches(1.2),
             "État de l'art & Marché", size=48, bold=True, color=WHITE)

    # ─── Slide 5 — Market cartography ───────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "ÉTAT DE L'ART", "Cartographie du marché — Outils PO & IA (2026)",
                "Quatre catégories d'outils, aucune ne couvre tout")
    cats = [
        ("ALM classiques", "Jira · Azure DevOps · Linear", "❌ Pas d'IA décisionnelle native"),
        ("GenAI intégrée", "Jira AI · Notion AI · ClickUp Brain", "❌ Génératif passif, pas d'agents"),
        ("PO-specific AI", "ProductBoard · Aha! · Productlogic", "❌ Black-box, SaaS US coûteux"),
        ("Frameworks agentiques", "LangChain · LangGraph · CrewAI", "❌ Génériques, pas métier PO"),
    ]
    for i, (h, ex, lim) in enumerate(cats):
        col = i % 2
        row = i // 2
        left = Inches(0.6 + col * 6.2)
        top = Inches(2.6 + row * 1.9)
        add_rect(s, left, top, Inches(6.0), Inches(1.65), PANEL)
        add_text(s, left + Inches(0.25), top + Inches(0.15), Inches(5.5), Inches(0.4),
                 h, size=15, bold=True, color=CYAN)
        add_text(s, left + Inches(0.25), top + Inches(0.6), Inches(5.5), Inches(0.4),
                 ex, size=12, color=WHITE)
        add_text(s, left + Inches(0.25), top + Inches(1.05), Inches(5.5), Inches(0.4),
                 lim, size=11, color=ORANGE)
    add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
             "→ Aucun outil ne combine RAG métier + multi-algorithmes transparents + vrais agents autonomes.",
             size=12, bold=True, color=CYAN)
    add_footer(s, prs, FOOTER)

    # ─── Slide 6 — GenAI integrated ─────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "ÉTAT DE L'ART", "GenAI intégrée dans les outils PO",
                "Atlassian Intelligence · Notion AI · Linear AI")
    rows = [
        ("Atlassian Intelligence (Jira AI)",
         "✅ Résumé d'issues, génération de descriptions, suggestions de story points",
         "❌ Pas de scoring multi-algos, pas d'agents, pas de RAG sur docs internes"),
        ("Notion AI / ClickUp Brain",
         "✅ Q&A sur workspace, génération de tâches, résumés",
         "❌ RAG limité, pas de priorisation quantitative"),
        ("Linear AI",
         "✅ Triage automatique, suggestion de cycles, UX excellente",
         "❌ Propriétaire, non extensible, hébergement US, opaque"),
    ]
    for i, (name, pro, con) in enumerate(rows):
        top = Inches(2.6 + i * 1.35)
        add_rect(s, Inches(0.6), top, Inches(12.1), Inches(1.2), PANEL)
        add_text(s, Inches(0.8), top + Inches(0.1), Inches(11.5), Inches(0.4),
                 name, size=14, bold=True, color=CYAN)
        add_text(s, Inches(0.8), top + Inches(0.55), Inches(11.5), Inches(0.3),
                 pro, size=11, color=WHITE)
        add_text(s, Inches(0.8), top + Inches(0.85), Inches(11.5), Inches(0.3),
                 con, size=11, color=ORANGE)
    add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
             "→ La GenAI intégrée reste une assistante passive : elle ne MODIFIE pas le backlog.",
             size=12, bold=True, color=CYAN)
    add_footer(s, prs, FOOTER)

    # ─── Slide 7 — PO-specific ──────────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "ÉTAT DE L'ART", "Outils PO-spécifiques",
                "ProductBoard · Aha! · Productlogic")
    rows = [
        ("ProductBoard", "Scoring : RICE",
         "✅ Backlog · Jira · roadmap",
         "❌ Black-box · SaaS US · 200–600 €/user/mois"),
        ("Aha!", "Scoring : custom",
         "✅ Roadmap AI · génération d'epics",
         "❌ Non transparent · très coûteux · hébergement US"),
        ("Productlogic", "Scoring : RICE + ICE",
         "✅ Multi-méthodes · interface claire",
         "❌ Pas de RAG · pas d'agents · pas d'XAI"),
    ]
    for i, (name, sc, pro, con) in enumerate(rows):
        left = Inches(0.6 + i * 4.15)
        add_rect(s, left, Inches(2.6), Inches(3.9), Inches(3.8), PANEL)
        add_text(s, left + Inches(0.25), Inches(2.8), Inches(3.5), Inches(0.4),
                 name, size=16, bold=True, color=CYAN)
        add_text(s, left + Inches(0.25), Inches(3.3), Inches(3.5), Inches(0.4),
                 sc, size=11, color=MUTED)
        add_text(s, left + Inches(0.25), Inches(3.9), Inches(3.5), Inches(1.2),
                 pro, size=11, color=WHITE)
        add_text(s, left + Inches(0.25), Inches(5.1), Inches(3.5), Inches(1.3),
                 con, size=11, color=ORANGE)
    add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             "→ Gap : aucun ne combine multi-algorithmes transparents + RAG documentaire + agents IA.",
             size=12, bold=True, color=CYAN)
    add_footer(s, prs, FOOTER)

    # ─── Slide 8 — RAG state of the art ─────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "ÉTAT DE L'ART", "RAG en 2026", "Évolution et choix techniques de PO.ai")
    add_rect(s, Inches(0.6), Inches(2.6), Inches(6.0), Inches(3.8), PANEL)
    add_text(s, Inches(0.8), Inches(2.8), Inches(5.6), Inches(0.4),
             "Évolution du RAG", size=14, bold=True, color=CYAN)
    add_text(s, Inches(0.8), Inches(3.3), Inches(5.6), Inches(3.0),
             "• Naive RAG → Advanced RAG → Agentic RAG\n"
             "• Re-ranking (Cohere, ColBERT)\n"
             "• MMR pour diversité du contexte\n"
             "• Hybrid search (BM25 + dense)\n"
             "• GraphRAG (Microsoft, 2024)\n"
             "• Embeddings : text-embedding-3-small (1536 d),\n   BGE-M3, Mistral-embed",
             size=12, color=WHITE)
    add_rect(s, Inches(6.85), Inches(2.6), Inches(6.0), Inches(3.8), PANEL)
    add_text(s, Inches(7.05), Inches(2.8), Inches(5.6), Inches(0.4),
             "Choix PO.ai", size=14, bold=True, color=VIOLET)
    add_text(s, Inches(7.05), Inches(3.3), Inches(5.6), Inches(3.0),
             "• ChromaDB (open source, persistant disque)\n"
             "• Embeddings : OpenAI text-embedding-3-small\n"
             "• Chunks 800 tokens · overlap 120\n"
             "• Retrieval MMR k=6 (diversité)\n"
             "• Routing LLM auto : GPT-4o / Mistral Large\n"
             "• Sources citées dans l'UI (chips cliquables)",
             size=12, color=WHITE)
    add_footer(s, prs, FOOTER)

    # ─── Slide 9 — OVH ──────────────────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "ÉTAT DE L'ART", "OVHcloud — Contexte du stage IAWF",
                "Outils PO existants et écosystème IA souverain")
    add_rect(s, Inches(0.6), Inches(2.6), Inches(6.0), Inches(3.8), PANEL)
    add_text(s, Inches(0.8), Inches(2.8), Inches(5.6), Inches(0.4),
             "Outils PO disponibles aujourd'hui", size=14, bold=True, color=ORANGE)
    add_text(s, Inches(0.8), Inches(3.3), Inches(5.6), Inches(3.0),
             "• Confluence  (documentation)\n"
             "• Jira  (tickets et backlog)\n\n"
             "❌ Aucun outil d'aide à la décision\n"
             "❌ Aucun RAG documentaire\n"
             "❌ Aucune priorisation quantitative\n\n"
             "→ Le PO est livré à son intuition.",
             size=12, color=WHITE)
    add_rect(s, Inches(6.85), Inches(2.6), Inches(6.0), Inches(3.8), OVH)
    add_text(s, Inches(7.05), Inches(2.8), Inches(5.6), Inches(0.4),
             "OVHcloud AI Endpoints (GA 2025)", size=14, bold=True, color=WHITE)
    add_text(s, Inches(7.05), Inches(3.3), Inches(5.6), Inches(3.0),
             "• API compatibles OpenAI\n"
             "• Mistral · Llama 3 · Mixtral · Qwen\n"
             "• Hébergement européen souverain\n"
             "  (Gravelines · Roubaix)\n"
             "• AI Notebooks / AI Training / AI Deploy\n"
             "• Managed DB + pgvector\n\n"
             "→ Migration PO.ai possible sans réécriture.",
             size=12, color=WHITE)
    add_footer(s, prs, FOOTER)

    # ─── Slide 10 — Positioning matrix ──────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "POSITIONNEMENT", "Ce que PO.ai fait — et que personne ne fait", "")
    headers = ["Dimension", "ProductBoard", "Aha!", "Productlogic", "PO.ai"]
    rows = [
        ("Algos transparents", "1 (RICE)", "Custom", "2", "7 ✅"),
        ("Agents IA (ReAct/LangGraph)", "❌", "❌", "❌", "✅"),
        ("RAG documentaire interne", "❌", "❌", "❌", "✅"),
        ("Sprint Planner (knapsack)", "❌", "❌", "❌", "✅"),
        ("Souveraineté (OVH compatible)", "❌", "❌", "❌", "✅"),
        ("Coût mensuel", "200–600€/u", "200–700€/u", "100–300€/u", "Self-hosted"),
    ]
    col_w = [Inches(3.4), Inches(2.0), Inches(1.8), Inches(2.2), Inches(2.8)]
    x0 = Inches(0.6)
    y0 = Inches(2.5)
    # Header
    x = x0
    for j, h in enumerate(headers):
        add_rect(s, x, y0, col_w[j], Inches(0.45), VIOLET)
        add_text(s, x + Inches(0.1), y0 + Inches(0.08), col_w[j] - Inches(0.1), Inches(0.3),
                 h, size=11, bold=True, color=WHITE)
        x += col_w[j]
    # Rows
    for i, row in enumerate(rows):
        y = y0 + Inches(0.5 + i * 0.55)
        x = x0
        for j, cell in enumerate(row):
            fill = PANEL if i % 2 == 0 else BG
            add_rect(s, x, y, col_w[j], Inches(0.5), fill)
            col = CYAN if j == 4 else WHITE
            add_text(s, x + Inches(0.1), y + Inches(0.1), col_w[j] - Inches(0.1), Inches(0.3),
                     cell, size=11, bold=(j == 4), color=col)
            x += col_w[j]
    add_footer(s, prs, FOOTER)

    # ─── Slide 11 — Part 2 ──────────────────────────────────────────────
    s = add_slide(prs)
    add_text(s, Inches(0.6), Inches(2.8), Inches(12), Inches(0.5),
             "PARTIE 2", size=18, bold=True, color=CYAN)
    add_text(s, Inches(0.6), Inches(3.3), Inches(12), Inches(1.2),
             "Architecture & Stack technique", size=44, bold=True, color=WHITE)

    # ─── Slide 12 — Architecture layers ─────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "ARCHITECTURE", "3 couches découplées",
                "Frontend Next.js  ·  Backend FastAPI  ·  Data & LLM")
    layers = [
        ("FRONTEND", "Next.js 15  ·  React 19  ·  TypeScript strict  ·  Tailwind CSS  ·  Lucide-react",
         "Chat · Prioritization · Sprint Planner · Retro · Documents", VIOLET),
        ("BACKEND", "FastAPI 0.115  ·  Pydantic 2.9  ·  LangChain 0.3  ·  LangGraph 0.5",
         "Routes : /api/chat · /api/prioritization · /api/documents · /api/agents", CYAN),
        ("DATA & LLM", "ChromaDB 0.5  ·  backlog.json  ·  OpenAI GPT-4o  ·  Mistral Large  ·  text-embedding-3-small",
         "Vector store local · Modèles routés selon la tâche · Migrable vers OVH AI Endpoints", OVH),
    ]
    for i, (h, stack, sub, color) in enumerate(layers):
        top = Inches(2.5 + i * 1.45)
        add_rect(s, Inches(0.6), top, Inches(12.1), Inches(1.3), PANEL)
        add_rect(s, Inches(0.6), top, Inches(0.18), Inches(1.3), color)
        add_text(s, Inches(1.0), top + Inches(0.12), Inches(11.5), Inches(0.4),
                 h, size=14, bold=True, color=color)
        add_text(s, Inches(1.0), top + Inches(0.55), Inches(11.5), Inches(0.4),
                 stack, size=12, color=WHITE)
        add_text(s, Inches(1.0), top + Inches(0.95), Inches(11.5), Inches(0.3),
                 sub, size=11, color=MUTED)
    add_text(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.3),
             "Principe : aucun couplage UI ↔ logique métier. Tous les modules testables indépendamment.",
             size=10, color=MUTED)

    # ─── Slide 13 — Tech stack table ────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "ARCHITECTURE", "Stack technique — Choix justifiés", "")
    headers = ["Composant", "Choix retenu", "Justification"]
    rows = [
        ("UI", "Next.js 15 + React 19 + Tailwind", "App Router moderne, TypeScript strict, dark mode natif"),
        ("API", "FastAPI 0.115 + Pydantic 2.9", "Async, validation stricte, auto-docs OpenAPI"),
        ("Orchestration LLM", "LangChain 0.3 + LangGraph 0.5", "Standard mature, state machines pour agents"),
        ("LLM raisonnement", "GPT-4o (OpenAI)", "Excellent pour RAG long-form et grooming"),
        ("LLM structuré", "Mistral Large", "JSON fiable · disponible sur OVH AI Endpoints"),
        ("Embeddings", "text-embedding-3-small", "1536 d, excellent ratio qualité/coût"),
        ("Vector store", "ChromaDB 0.5", "Open source, persistant, migrable pgvector"),
        ("ML scoring", "scikit-learn 1.8 (GradientBoosting)", "Algorithme ML Hybrid sur historique"),
        ("Secrets", "python-dotenv (.env)", "12-Factor App, jamais de clé dans le code"),
    ]
    col_w = [Inches(2.6), Inches(4.2), Inches(5.4)]
    x0 = Inches(0.6)
    y0 = Inches(2.3)
    x = x0
    for j, h in enumerate(headers):
        add_rect(s, x, y0, col_w[j], Inches(0.4), VIOLET)
        add_text(s, x + Inches(0.1), y0 + Inches(0.06), col_w[j] - Inches(0.1), Inches(0.3),
                 h, size=11, bold=True, color=WHITE)
        x += col_w[j]
    for i, row in enumerate(rows):
        y = y0 + Inches(0.45 + i * 0.45)
        x = x0
        for j, cell in enumerate(row):
            fill = PANEL if i % 2 == 0 else BG
            add_rect(s, x, y, col_w[j], Inches(0.4), fill)
            add_text(s, x + Inches(0.1), y + Inches(0.08), col_w[j] - Inches(0.1), Inches(0.3),
                     cell, size=10, color=WHITE, bold=(j == 1))
            x += col_w[j]
    add_footer(s, prs, FOOTER)

    # ─── Slide 14 — Part 3 ──────────────────────────────────────────────
    s = add_slide(prs)
    add_text(s, Inches(0.6), Inches(2.8), Inches(12), Inches(0.5),
             "PARTIE 3", size=18, bold=True, color=CYAN)
    add_text(s, Inches(0.6), Inches(3.3), Inches(12), Inches(1.2),
             "Réalisation — Modules, Algorithmes & Agents", size=38, bold=True, color=WHITE)

    # ─── Slide 15 — RAG module ──────────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "MODULE RAG", "rag/engine.py — Pipeline complet",
                "Ingestion · Retrieval MMR · Génération · Citations")
    steps = [
        ("1", "Ingestion", "PDF · TXT · MD\nUTF-8 autodetect"),
        ("2", "Chunking", "800 tokens\noverlap 120"),
        ("3", "Embeddings", "text-embedding-3-small\n1536 dimensions"),
        ("4", "ChromaDB", "Vector store\npersistant disque"),
        ("5", "Retrieval", "MMR k=6\n(diversité)"),
        ("6", "Génération", "GPT-4o / Mistral\nrouting auto"),
    ]
    w = Inches(2.0)
    for i, (n, h, b) in enumerate(steps):
        left = Inches(0.5 + i * 2.1)
        add_rect(s, left, Inches(2.7), w, Inches(2.0), PANEL)
        add_text(s, left + Inches(0.2), Inches(2.85), w - Inches(0.2), Inches(0.4),
                 n, size=14, bold=True, color=CYAN)
        add_text(s, left + Inches(0.2), Inches(3.25), w - Inches(0.2), Inches(0.4),
                 h, size=13, bold=True, color=WHITE)
        add_text(s, left + Inches(0.2), Inches(3.8), w - Inches(0.2), Inches(1.0),
                 b, size=10, color=MUTED)
    add_rect(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.6), PANEL)
    add_text(s, Inches(0.85), Inches(5.25), Inches(11.5), Inches(0.4),
             "Innovations", size=13, bold=True, color=VIOLET)
    add_text(s, Inches(0.85), Inches(5.7), Inches(11.5), Inches(1.0),
             "• Routing automatique du modèle (route_model) selon la nature de la tâche\n"
             "• MMR k=6 pour maximiser la diversité du contexte récupéré\n"
             "• Sources renvoyées au frontend → chips cliquables (traçabilité)",
             size=11, color=WHITE)
    add_footer(s, prs, FOOTER)

    # ─── Slide 16 — 7 algorithms ────────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "PRIORISATION", "Les 7 algorithmes — tous normalisés 0–100", "")
    algos = [
        ("RICE v2", "(Reach × Impact × Confidence²) / log(Effort)",
         "Standard PO · pondération logarithmique de l'effort"),
        ("WSJF v2", "(BV + TC + RR) / JobSize × time_decay",
         "SAFe · décroissance du Cost of Delay dans le temps"),
        ("ICE", "Impact × Confidence² × Ease",
         "Version simplifiée, idéale pour les MVPs et startups"),
        ("Kano", "Must-be / Performance / Delighter",
         "Satisfaction client via gain & risque de dissatisfaction"),
        ("Value/Effort", "matrice 2×2 dynamique",
         "Identification visuelle des Quick Wins et Big Bets"),
        ("AI Blend v2", "GPT-4o + Mistral ensemble (CoT)",
         "Score consensuel via raisonnement chaîné des deux LLM"),
        ("ML Hybrid", "GradientBoosting(scores des 6 autres)",
         "Apprentissage des patterns de priorisation historiques"),
    ]
    for i, (name, formula, body) in enumerate(algos):
        col = i % 2
        row = i // 2
        left = Inches(0.6 + col * 6.2)
        top = Inches(2.4 + row * 1.05)
        add_rect(s, left, top, Inches(6.0), Inches(0.95), PANEL)
        add_text(s, left + Inches(0.2), top + Inches(0.08), Inches(2.0), Inches(0.4),
                 name, size=13, bold=True, color=CYAN)
        add_text(s, left + Inches(2.2), top + Inches(0.1), Inches(3.7), Inches(0.4),
                 formula, size=10, color=VIOLET)
        add_text(s, left + Inches(0.2), top + Inches(0.5), Inches(5.7), Inches(0.4),
                 body, size=10, color=WHITE)
    add_footer(s, prs, FOOTER)

    # ─── Slide 17 — Sprint Planner ──────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "ALGORITHME PHARE", "Sprint Planner — Problème du Sac à Dos 0/1",
                "Maximiser le score total sous contrainte de capacité sprint")
    add_rect(s, Inches(0.6), Inches(2.5), Inches(6.0), Inches(4.0), PANEL)
    add_text(s, Inches(0.8), Inches(2.7), Inches(5.6), Inches(0.4),
             "Formulation mathématique", size=13, bold=True, color=CYAN)
    add_text(s, Inches(0.8), Inches(3.2), Inches(5.6), Inches(3.0),
             "Pour n features, capacité C (story points) :\n\n"
             "DP[i][c] = max(\n"
             "   DP[i-1][c],                  (skip)\n"
             "   DP[i-1][c - wᵢ] + vᵢ    (take)\n"
             ")\n\n"
             "• wᵢ = effort feature i (story points)\n"
             "• vᵢ = score priorisé (RICE / WSJF / …)\n"
             "• Complexité : O(n × C) — solution exacte\n"
             "• Fallback glouton si n × C > 10⁶",
             size=11, color=WHITE)
    add_rect(s, Inches(6.85), Inches(2.5), Inches(6.0), Inches(4.0), PANEL)
    add_text(s, Inches(7.05), Inches(2.7), Inches(5.6), Inches(0.4),
             "Sortie SprintPlan", size=13, bold=True, color=VIOLET)
    add_text(s, Inches(7.05), Inches(3.2), Inches(5.6), Inches(3.0),
             "• selected : features choisies\n"
             "• deferred : features reportées\n"
             "• utilization : % capacité utilisée\n"
             "• total_score : score cumulé\n"
             "• total_effort : effort total sélectionné\n\n"
             "Bénéfice PO : décision optimale prouvable,\n"
             "justifiable face au management.",
             size=11, color=WHITE)
    add_footer(s, prs, FOOTER)

    # ─── Slide 18 — Agent classification ────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "AGENTS — HONNÊTETÉ ACADÉMIQUE", "Tous les 'agents' ne sont pas des agents",
                "Référence : Anthropic 'Building Effective Agents' (décembre 2024)")
    rows = [
        ("Retro Analyzer", "One-shot prompt",
         "Un seul appel LLM avec sortie structurée. Pas de boucle, pas d'outils.", ORANGE),
        ("Sprint Planner", "Algorithme pur",
         "Programmation dynamique (knapsack 0/1). Aucun LLM.", ORANGE),
        ("Tool-calling Assistant", "VRAI AGENT (ReAct)",
         "Boucle LLM ↔ tools, max 5 itérations. 3 outils exposés.", GREEN),
        ("Epic Grooming", "VRAI AGENT (LangGraph)",
         "State machine avec validation INVEST + self-correction (max 2 retries).", GREEN),
    ]
    for i, (name, kind, desc, color) in enumerate(rows):
        top = Inches(2.5 + i * 1.05)
        add_rect(s, Inches(0.6), top, Inches(12.1), Inches(0.95), PANEL)
        add_rect(s, Inches(0.6), top, Inches(0.18), Inches(0.95), color)
        add_text(s, Inches(1.0), top + Inches(0.1), Inches(3.5), Inches(0.4),
                 name, size=13, bold=True, color=WHITE)
        add_text(s, Inches(4.5), top + Inches(0.1), Inches(3.5), Inches(0.4),
                 kind, size=12, bold=True, color=color)
        add_text(s, Inches(1.0), top + Inches(0.5), Inches(11.5), Inches(0.4),
                 desc, size=11, color=MUTED)
    add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
             "→ Le mécanisme le plus simple qui résout le problème (Anthropic 2024).",
             size=11, color=CYAN)
    add_footer(s, prs, FOOTER)

    # ─── Slide 19 — Tool-calling ReAct ──────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "AGENT 1", "Tool-calling Assistant — Boucle ReAct",
                "agents/assistant.py · bind_tools([add_feature, re_score, search_docs])")
    add_rect(s, Inches(0.6), Inches(2.5), Inches(6.0), Inches(4.0), PANEL)
    add_text(s, Inches(0.8), Inches(2.7), Inches(5.6), Inches(0.4),
             "Boucle de raisonnement", size=13, bold=True, color=CYAN)
    add_text(s, Inches(0.8), Inches(3.2), Inches(5.6), Inches(3.0),
             "User message\n"
             "    ↓\n"
             "LLM (GPT-4o + bind_tools) — décide\n"
             "    ↓ tool call ?\n"
             "    ├─ oui → execute → observation ──┐\n"
             "    │                                  │\n"
             "    └─ non → réponse finale            │\n"
             "                                       │\n"
             "    ← (reboucle, max 5 itérations) ───┘",
             size=11, color=WHITE)
    add_rect(s, Inches(6.85), Inches(2.5), Inches(6.0), Inches(4.0), PANEL)
    add_text(s, Inches(7.05), Inches(2.7), Inches(5.6), Inches(0.4),
             "3 outils exposés (@tool LangChain)", size=13, bold=True, color=VIOLET)
    add_text(s, Inches(7.05), Inches(3.2), Inches(5.6), Inches(3.0),
             "• add_feature(name, reach, impact, …)\n"
             "    → ajoute une feature au backlog\n\n"
             "• re_score(feature_id, algo)\n"
             "    → recalcule avec un autre algorithme\n\n"
             "• search_docs(query)\n"
             "    → recherche RAG dans les documents\n\n"
             "Garde-fous : max 5 itérations · validation Pydantic stricte",
             size=11, color=WHITE)
    add_footer(s, prs, FOOTER)

    # ─── Slide 20 — Epic Grooming LangGraph ─────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "AGENT 2 — LE PLUS AVANCÉ", "Epic Grooming — State Machine LangGraph",
                "agents/grooming.py · 4 nœuds + conditional edge avec self-correction")
    add_rect(s, Inches(0.6), Inches(2.5), Inches(6.5), Inches(4.0), PANEL)
    add_text(s, Inches(0.8), Inches(2.7), Inches(6.0), Inches(0.4),
             "Graph LangGraph", size=13, bold=True, color=CYAN)
    add_text(s, Inches(0.8), Inches(3.2), Inches(6.0), Inches(3.2),
             "[START]\n"
             "   ↓\n"
             "[split]      découpe l'epic en user stories\n"
             "   ↓\n"
             "[criteria]   génère les AC (INVEST)  ←──┐\n"
             "   ↓                                     │\n"
             "[validate]   vérifie INVEST + AC          │\n"
             "   ↓ conditional edge                     │\n"
             "   ├─ 'ok'    → [rice]                    │\n"
             "   └─ 'retry' (retry < 2) ────────────────┘\n"
             "   ↓\n"
             "[rice]       score chaque story\n"
             "   ↓\n"
             "[END]   stories + AC + scores + trace",
             size=10, color=WHITE)
    add_rect(s, Inches(7.35), Inches(2.5), Inches(5.5), Inches(4.0), PANEL)
    add_text(s, Inches(7.55), Inches(2.7), Inches(5.1), Inches(0.4),
             "État partagé (TypedDict)", size=13, bold=True, color=VIOLET)
    add_text(s, Inches(7.55), Inches(3.2), Inches(5.1), Inches(3.2),
             "class _State(TypedDict):\n"
             "  epic_title: str\n"
             "  stories: list[Story]\n"
             "  criteria: list[list[str]]\n"
             "  scores: list[float]\n"
             "  feedback: str | None\n"
             "  retry_count: int\n"
             "  trace: list[AgentStep]\n\n"
             "→ trace remontée dans l'UI :\n"
             "   transparence + débuggabilité.",
             size=10, color=WHITE)
    add_footer(s, prs, FOOTER)

    # ─── Slide 21 — Frontend ────────────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "FRONTEND", "5 panneaux dans une UX cohérente",
                "Next.js 15 · React 19 · Tailwind · Lucide-react")
    panels = [
        ("💬  Chat", "RAG + toggle Agent ON\n(tool-calling visible)"),
        ("📊  Prioritization", "Backlog + pills algorithmes\n+ bouton Groom Epic"),
        ("📅  Sprint Planner", "Slider capacité\n+ knapsack visualisé"),
        ("🔁  Retro", "Forces · améliorations\n· actions concrètes"),
        ("📁  Documents", "Upload + indexation\n+ liste des docs"),
    ]
    for i, (h, b) in enumerate(panels):
        left = Inches(0.6 + i * 2.5)
        add_rect(s, left, Inches(2.7), Inches(2.3), Inches(3.0), PANEL)
        add_text(s, left + Inches(0.2), Inches(2.9), Inches(2.0), Inches(0.5),
                 h, size=14, bold=True, color=CYAN)
        add_text(s, left + Inches(0.2), Inches(3.7), Inches(2.0), Inches(1.8),
                 b, size=11, color=WHITE)
    add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
             "Design system : violet #7C3AED · dark mode · Inter · Lucide · 100% TypeScript strict.",
             size=11, color=MUTED)
    add_footer(s, prs, FOOTER)

    # ─── Slide 22 — Part 4 ──────────────────────────────────────────────
    s = add_slide(prs)
    add_text(s, Inches(0.6), Inches(2.8), Inches(12), Inches(0.5),
             "PARTIE 4", size=18, bold=True, color=CYAN)
    add_text(s, Inches(0.6), Inches(3.3), Inches(12), Inches(1.2),
             "Démo · Résultats · Conclusion", size=44, bold=True, color=WHITE)

    # ─── Slide 23 — Demo scenario ───────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "DÉMO LIVE — 5 MINUTES", "Scénario de démonstration", "")
    steps = [
        ("1", "Upload doc", "Upload d'un PDF SAFe → 42 chunks indexés en quelques secondes"),
        ("2", "Chat RAG", "« Différence WSJF vs RICE ? » → réponse avec sources citées"),
        ("3", "Priorisation", "Switch RICE → WSJF → backlog réorganisé instantanément"),
        ("4", "Sprint Planner", "Capacité 40 SP → knapsack sélectionne 7 features optimales"),
        ("5", "Agent ReAct", "« Ajoute Export CSV puis re-score » → tool calls visibles"),
        ("6", "Epic Grooming", "Epic SSO → 5 stories + AC + RICE, avec trace de l'agent"),
    ]
    for i, (n, h, b) in enumerate(steps):
        top = Inches(2.4 + i * 0.78)
        add_rect(s, Inches(0.6), top, Inches(12.1), Inches(0.68), PANEL)
        add_rect(s, Inches(0.6), top, Inches(0.7), Inches(0.68), VIOLET)
        add_text(s, Inches(0.75), top + Inches(0.18), Inches(0.5), Inches(0.4),
                 n, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.5), top + Inches(0.1), Inches(2.5), Inches(0.5),
                 h, size=12, bold=True, color=CYAN)
        add_text(s, Inches(4.1), top + Inches(0.18), Inches(8.5), Inches(0.4),
                 b, size=11, color=WHITE)
    add_footer(s, prs, FOOTER)

    # ─── Slide 24 — Results ─────────────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "RÉSULTATS", "Mesures observées", "")
    stats = [
        ("2.1 s", "Latence RAG p50", CYAN),
        ("< 200 ms", "Knapsack 100 features", CYAN),
        ("8–12 s", "Grooming complet", VIOLET),
        ("~0.012 €", "Coût / interaction", GREEN),
        ("7", "Algorithmes", VIOLET),
        ("4", "Modules agentiques", VIOLET),
        ("~3500", "LoC backend", MUTED),
        ("~2800", "LoC frontend", MUTED),
    ]
    for i, (v, lbl, col) in enumerate(stats):
        c = i % 4
        r = i // 4
        left = Inches(0.6 + c * 3.1)
        top = Inches(2.6 + r * 2.0)
        add_rect(s, left, top, Inches(2.95), Inches(1.85), PANEL)
        add_text(s, left + Inches(0.2), top + Inches(0.25), Inches(2.7), Inches(0.8),
                 v, size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.2), top + Inches(1.15), Inches(2.7), Inches(0.5),
                 lbl, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
             "Mesures sur laptop dev (Windows · Python 3.11 · GPT-4o + Mistral Large).",
             size=10, color=MUTED)
    add_footer(s, prs, FOOTER)

    # ─── Slide 25 — Limits ──────────────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "LIMITES", "Honnêteté académique",
                "Ce que PO.ai ne fait pas (encore)")
    limits = [
        ("Pas de multi-utilisateurs / SSO",
         "Une instance = un PO. Authentification non implémentée."),
        ("Pas de connecteur Jira / Confluence réel",
         "Upload manuel des documents. Connecteurs prévus en roadmap."),
        ("Persistance JSON",
         "backlog.json en local. Migration Postgres + pgvector planifiée."),
        ("Évaluation RAG non systématique",
         "Tests qualitatifs uniquement. RAGAS / TruLens prévus."),
        ("ML Hybrid sur données synthétiques",
         "GradientBoosting entraîné sur exemples générés, pas d'historique réel."),
        ("Couverture tests partielle",
         "Tests unitaires sur les algorithmes uniquement."),
    ]
    for i, (h, b) in enumerate(limits):
        col = i % 2
        row = i // 2
        left = Inches(0.6 + col * 6.2)
        top = Inches(2.5 + row * 1.4)
        add_rect(s, left, top, Inches(6.0), Inches(1.25), PANEL)
        add_text(s, left + Inches(0.2), top + Inches(0.15), Inches(5.6), Inches(0.4),
                 "⚠  " + h, size=12, bold=True, color=ORANGE)
        add_text(s, left + Inches(0.2), top + Inches(0.65), Inches(5.6), Inches(0.6),
                 b, size=11, color=WHITE)
    add_footer(s, prs, FOOTER)

    # ─── Slide 26 — Roadmap OVH ─────────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "ROADMAP", "Migration OVHcloud — 1 jour de travail",
                "Stack 100% portable vers le cloud souverain européen")
    cols = [
        ("Court terme", [
            "Connecteur Jira & Confluence (OAuth)",
            "Migration LLM → OVH AI Endpoints",
            "Évaluation RAG avec RAGAS",
        ], CYAN),
        ("Moyen terme", [
            "Multi-tenant (Postgres + pgvector)",
            "Apprentissage continu (feedback 👍/👎)",
            "Tests utilisateurs avec PO OVH",
        ], VIOLET),
        ("Long terme", [
            "Marketplace OVHcloud",
            "Mode multi-équipes",
            "Publication open-source",
        ], GREEN),
    ]
    for i, (h, items, color) in enumerate(cols):
        left = Inches(0.6 + i * 4.15)
        add_rect(s, left, Inches(2.5), Inches(3.9), Inches(3.4), PANEL)
        add_text(s, left + Inches(0.25), Inches(2.7), Inches(3.5), Inches(0.4),
                 h, size=14, bold=True, color=color)
        for j, it in enumerate(items):
            add_text(s, left + Inches(0.25), Inches(3.2 + j * 0.55), Inches(3.5), Inches(0.5),
                     "→  " + it, size=11, color=WHITE)
    add_rect(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.7), OVH)
    add_text(s, Inches(0.85), Inches(6.25), Inches(11.5), Inches(0.4),
             "Migration : base_url = https://oai.endpoints.kepler.ai.cloud.ovh.net/v1  +  Mistral  ·  "
             "Chroma → pgvector OVH  ·  Deploy : AI Deploy (Docker)",
             size=11, color=WHITE)
    add_footer(s, prs, FOOTER)

    # ─── Slide 27 — Contributions ───────────────────────────────────────
    s = add_slide(prs)
    title_slide(s, prs, "APPORTS DU MÉMOIRE", "Cinq contributions", "")
    contribs = [
        ("📊", "Comparatif systématique du marché 2026",
         "Analyse de 12+ outils PO/IA : ALM, GenAI intégrée, PO-specific, frameworks agentiques."),
        ("🔬", "Implémentation transparente de 7 algorithmes",
         "RICE v2, WSJF v2, ICE, Kano, Value/Effort, AI Blend, ML Hybrid — tous normalisés 0–100."),
        ("🤖", "Distinction académique workflow vs agent",
         "Classification honnête (Anthropic 2024) : 2 vrais agents (ReAct + LangGraph), 2 workflows."),
        ("🚀", "Architecture portable vers cloud souverain",
         "Stack 100% compatible OVH AI Endpoints. Migration en 1 jour via base_url."),
        ("✅", "Pertinence opérationnelle validée",
         "Confirmé par le tuteur OVHcloud : « C'est ce qui manque dans notre processus actuel »."),
    ]
    for i, (icon, h, b) in enumerate(contribs):
        top = Inches(2.4 + i * 0.85)
        add_rect(s, Inches(0.6), top, Inches(12.1), Inches(0.75), PANEL)
        add_text(s, Inches(0.85), top + Inches(0.18), Inches(0.7), Inches(0.5),
                 icon, size=20, color=WHITE)
        add_text(s, Inches(1.6), top + Inches(0.1), Inches(11.0), Inches(0.4),
                 h, size=13, bold=True, color=CYAN)
        add_text(s, Inches(1.6), top + Inches(0.45), Inches(11.0), Inches(0.4),
                 b, size=11, color=WHITE)
    add_footer(s, prs, FOOTER)

    # ─── Slide 28 — Closing ─────────────────────────────────────────────
    s = add_slide(prs)
    add_header_bar(s, prs)
    add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(1.2),
             "« Le meilleur agent est celui qu'on comprend. »",
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # KPI strip
    kpis = [("7", "Algorithmes"), ("4", "Modules agentiques"),
            ("2.1 s", "Latence RAG"), ("100%", "Open & souverain")]
    for i, (v, lbl) in enumerate(kpis):
        left = Inches(0.6 + i * 3.1)
        add_rect(s, left, Inches(3.0), Inches(2.95), Inches(1.6), PANEL)
        add_text(s, left + Inches(0.2), Inches(3.2), Inches(2.7), Inches(0.7),
                 v, size=28, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.2), Inches(4.0), Inches(2.7), Inches(0.5),
                 lbl, size=11, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.6),
             "Merci ! Questions ?", size=32, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
             "Aouichaoui Takwa · M2 MIAS Centrale Lille · OVHcloud IAWF · Juillet 2026",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
             "PO.ai — RAG · 7 algorithmes · ReAct · LangGraph · Knapsack",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)

    out = Path("PO_Copilot_Soutenance_FINAL.pptx")
    prs.save(out)
    print(f"✅ Saved: {out.resolve()}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
