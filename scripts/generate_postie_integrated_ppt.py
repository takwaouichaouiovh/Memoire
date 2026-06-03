from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


BG = RGBColor(15, 23, 42)
PANEL = RGBColor(30, 41, 59)
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(148, 163, 184)
ACCENT = RGBColor(34, 211, 238)
VIOLET = RGBColor(124, 58, 237)


def add_bg(slide, prs):
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()


def add_header(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.color.rgb = ACCENT


def add_title(slide, title, subtitle=""):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(12), Inches(1.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(38)
    r.font.bold = True
    r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(16)
        r2.font.color.rgb = MUTED


def add_panel(slide, left, top, width, height, color=PANEL):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_bullets(slide, bullets, left=0.95, top=2.0, width=11.4, height=4.7, size=22):
    add_panel(slide, Inches(left - 0.25), Inches(top - 0.2), Inches(width + 0.5), Inches(height + 0.3))
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = f"• {b}"
        r.font.size = Pt(size)
        r.font.color.rgb = WHITE


def add_two_columns(slide, left_title, left_lines, right_title, right_lines):
    add_panel(slide, Inches(0.7), Inches(2.0), Inches(6.1), Inches(4.9))
    add_panel(slide, Inches(6.95), Inches(2.0), Inches(5.7), Inches(4.9))

    lt = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.5), Inches(0.6))
    ltf = lt.text_frame
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = left_title
    lr.font.size = Pt(22)
    lr.font.bold = True
    lr.font.color.rgb = ACCENT

    rt = slide.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.1), Inches(0.6))
    rtf = rt.text_frame
    rp = rtf.paragraphs[0]
    rr = rp.add_run()
    rr.text = right_title
    rr.font.size = Pt(22)
    rr.font.bold = True
    rr.font.color.rgb = ACCENT

    lbox = slide.shapes.add_textbox(Inches(1.0), Inches(2.9), Inches(5.5), Inches(3.7))
    ltf2 = lbox.text_frame
    ltf2.word_wrap = True
    for i, line in enumerate(left_lines):
        p = ltf2.paragraphs[0] if i == 0 else ltf2.add_paragraph()
        r = p.add_run()
        r.text = f"• {line}"
        r.font.size = Pt(16)
        r.font.color.rgb = WHITE

    rbox = slide.shapes.add_textbox(Inches(7.2), Inches(2.9), Inches(5.1), Inches(3.7))
    rtf2 = rbox.text_frame
    rtf2.word_wrap = True
    for i, line in enumerate(right_lines):
        p = rtf2.paragraphs[0] if i == 0 else rtf2.add_paragraph()
        r = p.add_run()
        r.text = f"• {line}"
        r.font.size = Pt(16)
        r.font.color.rgb = WHITE


def add_footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(12), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED


def build(path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    head = "POSTIE — OVHcloud IAWF · M2 Management IA · Centrale Lille × Universite de Lille · 2025-2026"

    # 1 Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header(s, head)
    add_title(s, "POSTIE", "Un copilote IA pour Product Owners et Scrum Masters")
    add_bullets(
        s,
        [
            "RAG documentaire",
            "7 algorithmes de priorisation",
            "Agents IA (ReAct + LangGraph)",
        ],
        top=2.5,
        height=2.6,
        size=26,
    )
    add_footer(s, "Aouichaoui Takwa")

    # 2 Roadmap
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header(s, head)
    add_title(s, "Roadmap de presentation")
    add_bullets(
        s,
        [
            "01 Introduction — Contexte OVHcloud & problematique terrain",
            "02 Analyse & Objectifs — Besoins identifies et hypotheses",
            "03 Methodologie — DSR, architecture et sprints Agile",
            "04 Stack Technique — FastAPI, Next.js, LangChain, ChromaDB",
            "05 Developpement & Resultats — RAG, algorithmes, agents",
            "06 Conclusion — apports, limites et perspectives",
        ],
        size=18,
    )
    add_footer(s, "Structure 01 → 06")

    # 3 Intro story
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header(s, head)
    add_title(s, "01 Introduction", "Le probleme que POSTIE resout")
    add_bullets(
        s,
        [
            "POSTIE aide a repondre a une question essentielle : que developper en priorite avec peu de temps ?",
            "Dans une equipe produit, le vrai probleme n'est pas de trouver des idees.",
            "Le vrai probleme est de choisir les bonnes idees.",
            "POSTIE aide justement a faire ce choix de maniere structuree.",
        ],
        size=20,
    )
    add_footer(s, "Contexte terrain : decision produit sous contrainte")

    # 4 Three foundations
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header(s, head)
    add_title(s, "02 Bases construites", "Les trois fondations de POSTIE")
    add_bullets(
        s,
        [
            "Base 1 : backlog structure (effort, impact, portee, confiance)",
            "Base 2 : assistant conversationnel connecte aux documents (PDF, Word, Markdown)",
            "Base 3 : plusieurs methodes de priorisation (RICE, WSJF, ICE...) pour comparer les decisions",
        ],
        size=21,
    )
    add_footer(s, "Du besoin metier vers une decision explicable")

    # 5 Advanced features
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header(s, head)
    add_title(s, "03 Fonctionnalites avancees", "Ce qui a ete ajoute ensuite")
    add_bullets(
        s,
        [
            "Planificateur de sprint : avec une capacite ex. 20 points, selection du meilleur ensemble",
            "Agent de grooming : transforme un epic en user stories claires et testables",
            "Analyse de retrospective : extraction automatique des actions, risques, points positifs et blocages",
            "Assistant oriente action : peut modifier le backlog et relancer les calculs",
        ],
        size=18,
    )
    add_footer(s, "Passage d'un assistant passif a un assistant operationnel")

    # 6 Three major improvements
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header(s, head)
    add_title(s, "04 Ameliorations majeures", "Les 3 evolutions ajoutees aujourd'hui")
    add_bullets(
        s,
        [
            "Gestion des dependances : une fonctionnalite B ne peut plus etre choisie sans A si B depend de A",
            "Superviseur multi-agents : objective unique -> grooming, priorisation, planification, resume",
            "Integration Jira Issues : import des issues dans le backlog et export du backlog vers Jira",
        ],
        size=19,
    )
    add_footer(s, "Fiabilite + automatisation + integration terrain")

    # 7 Methodology
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header(s, head)
    add_title(s, "05 Methodologie", "Approche DSR et execution Agile")
    add_two_columns(
        s,
        "Design Science Research",
        [
            "Pertinence : adequation au besoin terrain",
            "Rigueur : formalisation et validations",
            "Iteration : amelioration continue",
        ],
        "Execution en sprints",
        [
            "6 sprints planifies et suivis",
            "Priorisation continue des livrables",
            "Taux de completion eleve",
        ],
    )
    add_footer(s, "Une demarche produit-tech pilotee par la valeur")

    # 8 Stack
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header(s, head)
    add_title(s, "Stack technique")
    add_two_columns(
        s,
        "Frontend",
        [
            "Next.js 15, React 19, TypeScript",
            "Interfaces: chat, priorisation, sprint, retro, documents",
            "UX orientee demonstration et usage metier",
        ],
        "Backend & IA",
        [
            "FastAPI, Python 3.11, LangChain, LangGraph",
            "ChromaDB + embeddings + GPT-4o/Mistral",
            "API agents, priorisation, docs, sessions",
        ],
    )
    add_footer(s, "Architecture separee frontend/backend pour scalabilite")

    # 9 Results
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header(s, head)
    add_title(s, "Resultats", "Developpement et effets observes")
    add_bullets(
        s,
        [
            "RAG operationnel avec sources citees",
            "Moteur multi-algorithmes de priorisation deploye",
            "Planification sous contraintes avec dependances",
            "Supervision multi-agents traçable",
            "Integrations externes prêtes pour usage terrain",
        ],
        size=19,
    )
    add_footer(s, "Un prototype devenu un assistant concret")

    # 10 Conclusion
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_header(s, head)
    add_title(s, "06 Conclusion")
    add_bullets(
        s,
        [
            "POSTIE est passe d'un prototype a un assistant concret et utilisable",
            "Il combine IA, optimisation, orchestration multi-agents et connexion aux outils du terrain",
            "Perspectives : industrialisation, evaluation systematique IA, extension des connecteurs",
        ],
        top=2.3,
        height=3.2,
        size=23,
    )
    add_footer(s, "Merci")

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


if __name__ == "__main__":
    output = Path(__file__).resolve().parents[1] / "POSTIE_presentation_integree.pptx"
    build(output)
    print(output)
