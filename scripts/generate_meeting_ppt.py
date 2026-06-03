from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


BG = RGBColor(15, 23, 42)
PANEL = RGBColor(30, 41, 59)
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(148, 163, 184)
ACCENT = RGBColor(34, 211, 238)


def add_bg(slide, prs):
    rect = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG
    rect.line.fill.background()


def add_title(slide, title: str, subtitle: str = ""):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.0), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(16)
        r2.font.color.rgb = MUTED


def add_bullets(slide, bullets: list[str]):
    panel = slide.shapes.add_shape(1, Inches(0.6), Inches(1.9), Inches(12.1), Inches(5.2))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = f"• {item}"
        r.font.size = Pt(22)
        r.font.color.rgb = WHITE


def add_footer(slide, text: str):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.color.rgb = ACCENT


def make_ppt(output: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_title(
        s,
        "Entretien Annuel - Synthese Missions & Perspectives",
        "POSTIE + Contribution IA Web Factory",
    )
    add_bullets(
        s,
        [
            "Perimetre: alternance IAWF + developpement de POSTIE",
            "Objectif: montrer les missions, realisations, competences et objectifs a venir",
            "Positionnement: profil hybride produit-tech orienté IA",
        ],
    )
    add_footer(s, "Aouichaoui Takwa - M2 MIAS - 2026")

    # Slide 2
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_title(s, "1) Missions de l'annee")
    add_bullets(
        s,
        [
            "Assister l'IA Web Factory sur les sujets produit et operationnels",
            "Participer a la gestion du backlog (cadrage, priorisation, suivi)",
            "Participer a la gestion de la documentation fonctionnelle et technique",
            "S'occuper de l'usage et de l'adaptation de la bibliotheque de composants",
            "Tester et contribuer a l'amelioration de l'outil IA IAWF",
            "En parallele: conception et developpement de POSTIE",
        ],
    )
    add_footer(s, "Missions definies pour l'annee")

    # Slide 3
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_title(s, "2) Bilan des realisations (periode ecoulee)")
    add_bullets(
        s,
        [
            "Contribution continue aux activites backlog/documentation de l'IAWF",
            "Participation aux tests et retours d'amelioration de l'outil IA IAWF",
            "POSTIE: application complete (chat RAG, priorisation, planification)",
            "POSTIE: nouvelles briques avancees (orchestration multi-agents)",
            "POSTIE: planification avec dependances + integration GitHub Issues",
            "Approche iterative avec validations regulieres (tests et stabilisation)",
        ],
    )
    add_footer(s, "Resultats concrets obtenus")

    # Slide 4
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_title(s, "3) Competences developpees")
    add_bullets(
        s,
        [
            "IA agentique: conception de workflows et orchestration d'agents",
            "PO skills: backlog management, priorisation, logique valeur/effort",
            "Travail en mode iteratif avec amelioration continue",
            "Coordination produit-tech (besoin metier, contraintes techniques, UX)",
            "Communication avec parties prenantes techniques et non techniques",
        ],
    )
    add_footer(s, "Evolution des competences sur la periode")

    # Slide 5
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_title(s, "4) Objectifs periode a venir (alternance + POSTIE)")
    add_bullets(
        s,
        [
            "Poursuivre l'accompagnement global des sujets IAWF (backlog, doc, qualite)",
            "Renforcer la robustesse, l'adoption et la lisibilite des outils IA",
            "Faire evoluer POSTIE vers une version encore plus fiable pour demos et usage reel",
            "Structurer davantage les retours utilisateurs pour piloter la roadmap",
            "Maintenir un delivery regulier avec priorites mesurables",
        ],
    )
    add_footer(s, "Objectifs fixes pour la prochaine periode")

    # Slide 6
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_title(s, "5) Objectifs de developpement des competences")
    add_bullets(
        s,
        [
            "Monter en competences fonctionnelles produit",
            "Maitriser l'evaluation systematique des systemes IA",
            "Developper le pilotage produit-tech (arbitrages, roadmap, alignement)",
            "Gagner en aisance en presentation executive (synthese, storytelling, Q&A)",
            "Consolider un profil d'interface entre metier, produit et IA",
        ],
    )
    add_footer(s, "Plan de progression pour la periode a venir")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[1] / "presentation_entretien_IAWF_POSTIE.pptx"
    make_ppt(out)
    print(str(out))
